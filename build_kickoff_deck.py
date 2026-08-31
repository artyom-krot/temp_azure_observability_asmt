#!/usr/bin/env python3
"""Build kickoff PowerPoint deck from template."""

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import copy, os

TEMPLATE = 'templates/Observability_ASMT - Kick-Off.pptx'
OUTPUT = '00-context/kickoff-deck.pptx'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

prs = Presentation(TEMPLATE)
ORIG = len(prs.slides)           # 15
layouts = prs.slide_layouts

# Save slide references before we start adding
orig_cover     = prs.slides[0]
orig_thanks    = prs.slides[14]

SW = int(prs.slide_width)        # 12192000
SH = int(prs.slide_height)       # 6858000
ML = 457200
CT = 1380000
CW = SW - 2 * ML                 # 11277600

# ── Colors ──────────────────────────────────────────────────
NAVY       = (0,  32,  96)
WHITE      = (255,255,255)
LGRAY      = (242,242,242)
DARK       = ( 51, 51, 51)   # always set explicitly — master defaults to white text
RED_C      = (180,  0,   0)
GREEN_C    = (  0,112,   0)
AMBER      = (191,144,  0)

# ── Utilities ───────────────────────────────────────────────

def clone(prs, src):
    """Clone src slide and append to prs."""
    new = prs.slides.add_slide(src.slide_layout)
    sp = new.shapes._spTree
    for ch in list(sp)[2:]:
        sp.remove(ch)
    for ch in list(src.shapes._spTree)[2:]:
        sp.append(copy.deepcopy(ch))
    for rId, rel in src.part.rels.items():
        if not rel.is_external:
            try:
                new.part.relate_to(rel.target_part, rel.reltype)
            except Exception:
                pass
    return new

def del_slide(prs, idx):
    lst = prs.slides._sldIdLst
    sldId = lst[idx]
    rId = sldId.get(f'{{{R_NS}}}id')
    lst.remove(sldId)
    prs.part.drop_rel(rId)

def sec(prs, num, title):
    """Section header slide."""
    s = prs.slides.add_slide(layouts[7])
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        # clear
        for p in tf.paragraphs[1:]:
            p._p.getparent().remove(p._p)
        tf.paragraphs[0].clear()
        run = tf.paragraphs[0].add_run()
        if 'Title' in sh.name:
            run.text = title
        elif 'Placeholder 2' in sh.name:
            run.text = num
    return s

def title_slide(prs, title_text):
    """Title-layout content slide."""
    s = prs.slides.add_slide(layouts[17])
    for sh in s.shapes:
        if sh.has_text_frame and sh.name.startswith('Title'):
            tf = sh.text_frame
            for p in tf.paragraphs[1:]:
                p._p.getparent().remove(p._p)
            tf.paragraphs[0].clear()
            tf.paragraphs[0].add_run().text = title_text
            break
    return s

def tb(s, l, t, w, h):
    return s.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))

def wp(tf, items, ww=True):
    """
    Write paragraphs. items = list of (text, dict) or just str.
    dict keys: size, bold, italic, color, bullet, sb (space_before), align
    NOTE: color defaults to DARK (51,51,51) — master theme defaults to white text
    which would be invisible on white slide backgrounds.
    """
    tf.word_wrap = ww
    first = True
    for item in items:
        text, kw = (item, {}) if isinstance(item, str) else item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        if kw.get('sb'):
            p.space_before = Pt(kw['sb'])
        if kw.get('align'):
            p.alignment = kw['align']
        if not text:
            continue
        r = p.add_run()
        r.text = ('• ' if kw.get('bullet') else '') + text
        r.font.size = Pt(kw.get('size', 14))
        r.font.bold  = kw.get('bold',  False)
        r.font.italic= kw.get('italic',False)
        # Always set explicit color — theme master defaults to white text
        r.font.color.rgb = RGBColor(*kw.get('color', DARK))

def tbl(s, rows, cols, l, t, w, h):
    return s.shapes.add_table(rows, cols, Emu(l), Emu(t), Emu(w), Emu(h)).table

def hdr(t, col, text, sz=11):
    c = t.cell(0, col)
    c.text_frame.word_wrap = True
    p = c.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = True
    r.font.color.rgb = RGBColor(*WHITE)   # explicit white on navy
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(*NAVY)

def cel(t, row, col, text, shade=False, sz=11, bold=False,
        align=PP_ALIGN.LEFT, bg=None, fg=None):
    c = t.cell(row, col)
    c.text_frame.word_wrap = True
    p = c.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    # Always set explicit text color — table style "Medium Style 2 - Accent 1"
    # has near-black (accent1=#060606) cell backgrounds; explicit fg overrides.
    r.font.color.rgb = RGBColor(*(fg if fg else DARK))
    # Always set explicit cell background — table style default is near-black.
    fill_color = bg if bg else (LGRAY if shade else WHITE)
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(*fill_color)


# ════════════════════════════════════════════════════════════
# SLIDE 1 — Cover  (cloned from template)
# ════════════════════════════════════════════════════════════
cover = clone(prs, orig_cover)
for sh in cover.shapes:
    if sh.has_text_frame and 'Placeholder 6' in sh.name:
        tf = sh.text_frame
        for p in tf.paragraphs[1:]:
            p._p.getparent().remove(p._p)
        tf.paragraphs[0].clear()
        tf.paragraphs[0].add_run().text = 'Aug 2026'

# ════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'Agenda')
box = tb(s, ML, CT, CW, 4800000)
agenda_items = [
    ('01', 'What We Know Right Now'),
    ('02', 'Problem Statement'),
    ('03', 'Scope Definition'),
    ('04', 'Expected Outcomes'),
    ('05', 'Assessment Plan and Approach'),
    ('06', 'RAID — Risks, Assumptions, Dependencies'),
]
wp(box.text_frame, [
    (f'{n}   {t}', {'size': 22, 'bold': False, 'sb': 8 if n != '01' else 0})
    for n, t in agenda_items
])

# ════════════════════════════════════════════════════════════
# SLIDE 3 — Section 01
# ════════════════════════════════════════════════════════════
sec(prs, '01', 'What We Know Right Now')

# ════════════════════════════════════════════════════════════
# SLIDE 4 — Platform Overview
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'Platform Overview')
t = tbl(s, 8, 3, ML, CT, CW, 4600000)
t.columns[0].width = Emu(2600000)
t.columns[1].width = Emu(5200000)
t.columns[2].width = Emu(3477600)
for i, h in enumerate(['Component', 'Role / Function', 'Status']):
    hdr(t, i, h, 12)
rows = [
    ('CDN', 'Edge layer — ingress, caching, DDoS protection', 'Provider TBC'),
    ('Apache', 'Web / reverse proxy tier', 'Active'),
    ('Azure VMs', 'Legacy application components', 'Migration → AKS in progress'),
    ('Azure Kubernetes Service', 'Modernised microservices, namespace-based env isolation', 'Active'),
    ('Azure SQL Database', 'Managed relational DB with HA', 'Active'),
    ('Azure Service Bus', 'Regional async messaging between services', 'Active'),
    ('App Gateway + WAF', 'L7 load balancing and web application firewall', 'Active'),
]
for ri, (comp, role, status) in enumerate(rows, 1):
    shade = ri % 2 == 0
    cel(t, ri, 0, comp, shade, 11, bold=True)
    cel(t, ri, 1, role, shade, 11)
    cel(t, ri, 2, status, shade, 11)
note = tb(s, ML, CT + 4700000, CW, 400000)
wp(note.text_frame, [
    ('12 production environments (primary + DR) across US, UK, AU  |  DE in progress, CA planned  |  Monitoring: Azure Monitor + Datadog single tenant',
     {'size': 10, 'italic': True, 'color': (100,100,100)})
])

# ════════════════════════════════════════════════════════════
# SLIDE 5 — Current Monitoring State
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'Current Monitoring State')
HW = int((CW - 120000) / 2)
L2 = ML + HW + 120000

lb = tb(s, ML, CT, HW, 4800000)
wp(lb.text_frame, [
    ('What is in place', {'size': 16, 'bold': True, 'color': GREEN_C}),
    ('', {}),
    ('Azure Monitor + Log Analytics for infrastructure signals',               {'size': 13, 'bullet': True}),
    ('Datadog as primary monitoring platform (migrated from Splunk ~2 yrs ago)',{'size': 13, 'bullet': True}),
    ('Single Datadog tenant covers ALL consumer environments',                  {'size': 13, 'bullet': True}),
    ('Tag-based environment isolation: env / consumer / service',               {'size': 13, 'bullet': True}),
    ('AKS namespace-based environment separation',                              {'size': 13, 'bullet': True}),
])

rb = tb(s, L2, CT, HW, 4800000)
wp(rb.text_frame, [
    ('What we do not know yet', {'size': 16, 'bold': True, 'color': RED_C}),
    ('', {}),
    ('Which tiers have no monitoring coverage (dark)',                          {'size': 13, 'bullet': True}),
    ('Whether all 12 environments are consistently configured',                 {'size': 13, 'bullet': True}),
    ('Actual alert coverage vs. real failure modes',                            {'size': 13, 'bullet': True}),
    ('Whether SLA obligations map to any observable signal',                    {'size': 13, 'bullet': True}),
    ('APM instrumentation depth and trace coverage',                            {'size': 13, 'bullet': True}),
    ('Incident history and MTTD / MTTR baseline',                              {'size': 13, 'bullet': True}),
    ('Whether tag governance is enforced or ad hoc',                           {'size': 13, 'bullet': True}),
])

# ════════════════════════════════════════════════════════════
# SLIDE 6 — Section 02
# ════════════════════════════════════════════════════════════
sec(prs, '02', 'Problem Statement')

# ════════════════════════════════════════════════════════════
# SLIDE 7 — Problem Statement
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'Problem Statement')
lb = tb(s, ML, CT, HW, 5000000)
wp(lb.text_frame, [
    ('Current State', {'size': 16, 'bold': True, 'color': RED_C}),
    ('', {}),
    ('Observability is reactive. Consumers report failures before internal monitoring detects them.',
     {'size': 13}),
    ('', {}),
    ('Confirmed evidence:', {'size': 13, 'bold': True}),
    ('Database crash caused consumer-facing outage — detected externally', {'size': 12, 'bullet': True}),
    ('', {}),
    ('Root causes (hypothesised — to validate):', {'size': 13, 'bold': True}),
    ('Instrumentation gaps — tiers emit insufficient signals',      {'size': 12, 'bullet': True}),
    ('Alert coverage gaps — critical failure modes not monitored',  {'size': 12, 'bullet': True}),
    ('Alert quality — noisy, mis-tuned, or misrouted',             {'size': 12, 'bullet': True}),
    ('SLA/SLO disconnect — no automated observability backing',    {'size': 12, 'bullet': True}),
    ('Governance gaps — no enforced monitoring standard',          {'size': 12, 'bullet': True}),
])
rb = tb(s, L2, CT, HW, 5000000)
wp(rb.text_frame, [
    ('Target State — North Star', {'size': 16, 'bold': True, 'color': GREEN_C}),
    ('', {}),
    ('Platform health continuously and proactively monitored across all tiers and all consumer environments',
     {'size': 13, 'bullet': True, 'sb': 4}),
    ('Every SLA obligation backed by an automated signal that alerts before breach occurs',
     {'size': 13, 'bullet': True, 'sb': 4}),
    ('Incidents detected internally — MTTD in minutes, not hours reported by consumers',
     {'size': 13, 'bullet': True, 'sb': 4}),
    ('Engineering, operations, and leadership each have purpose-built visibility',
     {'size': 13, 'bullet': True, 'sb': 4}),
    ('Observability standards governed, owned, and consistently applied across all teams',
     {'size': 13, 'bullet': True, 'sb': 4}),
])

# ════════════════════════════════════════════════════════════
# SLIDE 8 — Section 03
# ════════════════════════════════════════════════════════════
sec(prs, '03', 'Scope Definition')

# ════════════════════════════════════════════════════════════
# SLIDE 9 — In-Scope Activities
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'In-Scope Activities — Phase 1')
t = tbl(s, 9, 3, ML, CT, CW, 5000000)
t.columns[0].width = Emu(500000)
t.columns[1].width = Emu(4600000)
t.columns[2].width = Emu(6177600)
for i, h in enumerate(['#', 'Activity', 'What We Are Assessing']):
    hdr(t, i, h, 11)
acts = [
    ('1','Application & Infrastructure Landscape Review',
     'Identify all in-scope tiers, services, components — the factual baseline'),
    ('2','Observability Coverage Evaluation',
     'Log pipelines, metrics, tracing — identify dark tiers and instrumentation blind spots'),
    ('3','Alerting & Detection Gap Analysis',
     'Alert inventory vs. Critical/High failure modes; gaps where consumers detect before monitoring'),
    ('4','APM Review',
     'Evaluate whether APM depth supports effective root cause analysis'),
    ('5','SLA / SLO Observability Review',
     'Map SLA commitments to SLO monitors; identify gaps where obligations lack observable signals'),
    ('6','Dashboard & Visibility Assessment',
     'Evaluate dashboards against audience needs: engineering, operations, leadership'),
    ('7','Operational Processes Review',
     'Incident workflow: routing, on-call, runbooks, escalation paths, MTTD / MTTR posture'),
    ('8','Governance Model Review',
     'Current-state governance: tagging, ownership, access control, standards consistency'),
]
for ri, (num, act, focus) in enumerate(acts, 1):
    shade = ri % 2 == 0
    cel(t, ri, 0, num,   shade, 12, bold=True, align=PP_ALIGN.CENTER)
    cel(t, ri, 1, act,   shade, 11, bold=True)
    cel(t, ri, 2, focus, shade, 11)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — Scope Boundaries
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'Scope Boundaries and Constraints')
lb = tb(s, ML, CT, HW, 5000000)
wp(lb.text_frame, [
    ('In Scope', {'size': 16, 'bold': True, 'color': GREEN_C}),
    ('', {}),
    ('Platform tiers:', {'size': 13, 'bold': True}),
    ('CDN / Apache / AKS / Azure VMs',                    {'size': 12, 'bullet': True}),
    ('Azure SQL, Service Bus, Blob, App Gateway + WAF',   {'size': 12, 'bullet': True}),
    ('Datadog — monitors, APM, logs, SLOs',               {'size': 12, 'bullet': True}),
    ('Azure Monitor + Log Analytics',                      {'size': 12, 'bullet': True}),
    ('', {}),
    ('Active regions:', {'size': 13, 'bold': True}),
    ('US, UK, AU  (DE — to confirm today)',               {'size': 12, 'bullet': True}),
    ('', {}),
    ('Capacity:', {'size': 13, 'bold': True}),
    ('2 specialists  |  4 weeks of active delivery',      {'size': 12, 'bullet': True}),
])
rb = tb(s, L2, CT, HW, 5000000)
wp(rb.text_frame, [
    ('Out of Scope — Phase 2', {'size': 16, 'bold': True, 'color': (140,40,0)}),
    ('', {}),
    ('Multi-region consistency deep-dive',                     {'size': 13, 'bullet': True, 'sb': 4}),
    ('Datadog platform utilisation optimisation',              {'size': 13, 'bullet': True, 'sb': 4}),
    ('Monitoring-as-code / CI/CD automation',                  {'size': 13, 'bullet': True, 'sb': 4}),
    ('Business activity monitoring — consumer KPIs',           {'size': 13, 'bullet': True, 'sb': 4}),
    ('Unified observability architecture implementation',       {'size': 13, 'bullet': True, 'sb': 4}),
    ('Standardisation at scale / migration execution',         {'size': 13, 'bullet': True, 'sb': 4}),
])
note = tb(s, ML, CT + 5100000, CW, 350000)
wp(note.text_frame, [
    ('* Week 1 is a landscape discovery phase — if undisclosed components surface, both parties agree in writing whether to include them or defer to Phase 2.',
     {'size': 10, 'italic': True, 'color': (100,100,100)})
])

# ════════════════════════════════════════════════════════════
# SLIDE 11 — Section 04
# ════════════════════════════════════════════════════════════
sec(prs, '04', 'Expected Outcomes')

# ════════════════════════════════════════════════════════════
# SLIDE 12 — Deliverables
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'Expected Outcomes and Deliverables')
t = tbl(s, 6, 3, ML, CT, CW, 3400000)
t.columns[0].width = Emu(700000)
t.columns[1].width = Emu(5200000)
t.columns[2].width = Emu(5377600)
for i, h in enumerate(['#', 'Deliverable', 'Primary Audience']):
    hdr(t, i, h, 12)
delivs = [
    ('1','Technical Assessment Report',      'Engineering & Architecture'),
    ('2','Target Observability Architecture','Engineering & Architecture'),
    ('3','Improvement Roadmap',              'Engineering & Leadership'),
    ('4','Governance Guidelines',            'Engineering & Operations'),
    ('5','Executive Summary',               'Leadership & Stakeholders'),
]
for ri, (num, name, aud) in enumerate(delivs, 1):
    shade = ri % 2 == 0
    cel(t, ri, 0, num,  shade, 12, bold=True, align=PP_ALIGN.CENTER)
    cel(t, ri, 1, name, shade, 12, bold=True)
    cel(t, ri, 2, aud,  shade, 12)

# Two columns below table
BOT_T = CT + 3600000
lb2 = tb(s, ML, BOT_T, HW, 2400000)
wp(lb2.text_frame, [
    ('Technical gaps identified', {'size': 14, 'bold': True, 'color': NAVY}),
    ('Dark tiers with no monitoring coverage',       {'size': 12, 'bullet': True}),
    ('Missing or misconfigured alert rules',          {'size': 12, 'bullet': True}),
    ('SLA obligations without observable signals',    {'size': 12, 'bullet': True}),
    ('APM instrumentation and trace coverage gaps',   {'size': 12, 'bullet': True}),
])
rb2 = tb(s, L2, BOT_T, HW, 2400000)
wp(rb2.text_frame, [
    ('Process and governance gaps identified', {'size': 14, 'bold': True, 'color': NAVY}),
    ('Ownership and governance model gaps',         {'size': 12, 'bullet': True}),
    ('Runbook quality and on-call effectiveness',   {'size': 12, 'bullet': True}),
    ('Absence of SLO-driven alerting practice',     {'size': 12, 'bullet': True}),
    ('Quick wins for immediate improvement',        {'size': 12, 'bullet': True}),
])

# ════════════════════════════════════════════════════════════
# SLIDE 13 — Section 05
# ════════════════════════════════════════════════════════════
sec(prs, '05', 'Assessment Plan')

# ════════════════════════════════════════════════════════════
# SLIDE 14 — Week-by-week plan
# ════════════════════════════════════════════════════════════
s = title_slide(prs, '4-Week Assessment Plan')
t = tbl(s, 5, 3, ML, CT, CW, 4200000)
t.columns[0].width = Emu(1400000)
t.columns[1].width = Emu(2600000)
t.columns[2].width = Emu(7277600)
for i, h in enumerate(['Week', 'Focus', 'Key Activities']):
    hdr(t, i, h, 12)
plan = [
    ('Week 1','Discovery & Landscape',
     'Kick-off  ·  Engineering Workshop  ·  Access provisioning  ·  Platform landscape mapping  ·  Initial evidence collection'),
    ('Week 2','Technical Analysis',
     'Operations Workshop  ·  Datadog deep-dive (monitors, APM, logs, SLOs)  ·  Azure Monitor review  ·  AKS / source code / CI/CD review'),
    ('Week 3','Gap Analysis',
     'Structured gap analysis across all domains  ·  Findings drafting  ·  Interim Findings Review with Client'),
    ('Week 4','Recommendations & Delivery',
     'Target architecture  ·  Roadmap  ·  Report finalisation  ·  Leadership Workshop (conditional)  ·  Final Presentation'),
]
for ri, (wk, focus, acts) in enumerate(plan, 1):
    shade = ri % 2 == 0
    cel(t, ri, 0, wk,    shade, 12, bold=True)
    cel(t, ri, 1, focus, shade, 12, bold=True)
    cel(t, ri, 2, acts,  shade, 11)
note = tb(s, ML, CT + 4300000, CW, 380000)
wp(note.text_frame, [
    ('* Timeline is indicative. Client-side delays in Weeks 1–2 will be flagged by end of Week 2 with an updated delivery date.',
     {'size': 10, 'italic': True, 'color': (100,100,100)})
])

# ════════════════════════════════════════════════════════════
# SLIDE 15 — Workshops & Toolset
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'Workshops and Toolset')
TW = int((CW - 200000) * 0.46)
L2B = ML + TW + 200000
RW = CW - TW - 200000

wt = tbl(s, 7, 2, ML, CT, TW, 4200000)
wt.columns[0].width = Emu(int(TW * 0.65))
wt.columns[1].width = Emu(TW - int(TW * 0.65))
hdr(wt, 0, 'Workshop', 11); hdr(wt, 1, 'Timing', 11)
workshops = [
    ('Kick-off',                           'Week 1'),
    ('Engineering Workshop',               'Week 1'),
    ('Operations / On-call Workshop',      'Week 2'),
    ('Interim Findings Review',            'Week 3'),
    ('Leadership Workshop (conditional)',  'Wk 3–4'),
    ('Final Delivery Presentation',        'Week 4'),
]
for ri, (w, timing) in enumerate(workshops, 1):
    shade = ri % 2 == 0
    cel(wt, ri, 0, w,      shade, 11)
    cel(wt, ri, 1, timing, shade, 11, align=PP_ALIGN.CENTER)

rb = tb(s, L2B, CT, RW, 4800000)
wp(rb.text_frame, [
    ('Toolset', {'size': 16, 'bold': True, 'color': NAVY}),
    ('', {}),
    ('Azure CLI + Portal',     {'size': 13, 'bold': True}),
    ('Diagnostic settings, Azure Monitor, Log Analytics, AKS', {'size': 11, 'color': (80,80,80)}),
    ('Datadog API (read-only)', {'size': 13, 'bold': True, 'sb': 5}),
    ('Programmatic export of monitors, dashboards, SLOs, APM, log indexes', {'size': 11, 'color': (80,80,80)}),
    ('Datadog UI — Viewer role', {'size': 13, 'bold': True, 'sb': 5}),
    ('Direct inspection and screenshot evidence', {'size': 11, 'color': (80,80,80)}),
    ('Source code + CI/CD (read-only)', {'size': 13, 'bold': True, 'sb': 5}),
    ('Instrumentation verification; observability-as-code assessment', {'size': 11, 'color': (80,80,80)}),
    ('Claude — Anthropic AI', {'size': 13, 'bold': True, 'sb': 5}),
    ('Evidence analysis, gap identification, report drafting — client confirmation required', {'size': 11, 'color': (80,80,80)}),
])

# ════════════════════════════════════════════════════════════
# SLIDE 16 — Section 06
# ════════════════════════════════════════════════════════════
sec(prs, '06', 'RAID — Risks, Assumptions, Dependencies')

# ════════════════════════════════════════════════════════════
# SLIDE 17 — Risks
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'Risks')
t = tbl(s, 9, 4, ML, CT, CW, 5000000)
t.columns[0].width = Emu(450000)
t.columns[1].width = Emu(3800000)
t.columns[2].width = Emu(4800000)
t.columns[3].width = Emu(2227600)
for i, h in enumerate(['#', 'Risk', 'Impact', 'Owner']):
    hdr(t, i, h, 11)
risks = [
    ('R1','Azure / Datadog access not provisioned by Week 1 Day 1',
     'Timeline compresses or slips entirely', 'Client'),
    ('R2','Key workshop participants unavailable or substituted with proxies',
     'Findings based on assumptions; reduced accuracy', 'Client'),
    ('R3','SLA/SLO definitions, architecture docs, or incident records unavailable',
     'Gap analysis lacks contractual grounding; findings understate severity', 'Client'),
    ('R4','Platform scope larger than known — services discovered in Week 1',
     'Gaps in unknown components go unassessed', 'Both'),
    ('R5','IAM / security policies block required read-only access',
     'Domain coverage reduced; findings may be caveated', 'Client'),
    ('R6','Single Datadog tenant complexity makes per-environment analysis difficult',
     'Environment-specific gaps may go undetected', 'Both'),
    ('R7','AI tooling not cleared by Client data handling policy',
     'Methodology revised mid-delivery; timeline extends', 'Both'),
    ('R8','Delays in Weeks 1–2 compress analysis and delivery weeks',
     'Report and architecture quality reduced', 'Both'),
]
for ri, (num, risk, impact, owner) in enumerate(risks, 1):
    shade = ri % 2 == 0
    cel(t, ri, 0, num,    shade, 11, bold=True, align=PP_ALIGN.CENTER)
    cel(t, ri, 1, risk,   shade, 10)
    cel(t, ri, 2, impact, shade, 10)
    cel(t, ri, 3, owner,  shade, 11, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 18 — Dependencies
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'Assumptions and Dependencies')
t = tbl(s, 10, 3, ML, CT, CW, 5100000)
t.columns[0].width = Emu(5200000)
t.columns[1].width = Emu(2000000)
t.columns[2].width = Emu(4077600)
for i, h in enumerate(['Dependency', 'Blocking?', 'Required by']):
    hdr(t, i, h, 11)
deps = [
    ('Azure read-only access (Reader, Log Analytics, AKS Cluster User, Policy Reader)',
     'BLOCKING', 'Week 1, Day 1'),
    ('Datadog API Key + Application Key (read-only)', 'BLOCKING', 'Week 1, Day 1'),
    ('Datadog Viewer role (UI access)',               'BLOCKING', 'Week 1, Day 1'),
    ('Source code repository read access',            'BLOCKING', 'Week 1'),
    ('In-scope subscription IDs list',                'BLOCKING', 'Kick-off'),
    ('AI tooling acceptable-use confirmation',         'BLOCKING', 'Before Week 1'),
    ('Named Engineering Workshop participants',        'BLOCKING', 'Before Week 1 ends'),
    ('Named Operations Workshop participants',         'BLOCKING', 'Before Week 2'),
    ('SLA/SLO definitions + incident post-mortems',   'Required', 'Week 2'),
]
for ri, (dep, blocking, req) in enumerate(deps, 1):
    shade = ri % 2 == 0
    is_blocking = blocking == 'BLOCKING'
    bg_row = (255, 230, 230) if is_blocking else None
    cel(t, ri, 0, dep, shade, 10, bg=bg_row)
    cel(t, ri, 1, blocking, shade, 10, bold=is_blocking,
        fg=(180, 0, 0) if is_blocking else None,
        bg=bg_row, align=PP_ALIGN.CENTER)
    cel(t, ri, 2, req, shade, 10, bg=bg_row)

# ════════════════════════════════════════════════════════════
# SLIDE 19 — Next Steps
# ════════════════════════════════════════════════════════════
s = title_slide(prs, 'Next Steps — To Close in This Session')
box = tb(s, ML, CT, CW, 5200000)
wp(box.text_frame, [
    ('Actions required before Week 1 begins:', {'size': 16, 'bold': True, 'color': NAVY}),
    ('', {}),
    ('Client confirms or adjusts platform scope — tiers, regions, environments',
     {'size': 16, 'bullet': True, 'sb': 4}),
    ('Client confirms AI tooling acceptable-use policy',
     {'size': 16, 'bullet': True, 'sb': 4}),
    ('Client nominates named participants for Engineering and Operations workshops',
     {'size': 16, 'bullet': True, 'sb': 4}),
    ('Client names a single point of contact for access provisioning',
     {'size': 16, 'bullet': True, 'sb': 4}),
    ('Client provides list of in-scope Azure subscription IDs',
     {'size': 16, 'bullet': True, 'sb': 4}),
    ('Both parties agree access provisioning deadline — target: 2 business days post kick-off',
     {'size': 16, 'bullet': True, 'sb': 4}),
    ('Both parties confirm date and time for Week 1 Engineering Workshop',
     {'size': 16, 'bullet': True, 'sb': 4}),
    ('Client confirms documentation owner for SLA/SLO definitions and incident records',
     {'size': 16, 'bullet': True, 'sb': 4}),
])

# ════════════════════════════════════════════════════════════
# SLIDE 20 — Thank You  (cloned from template)
# ════════════════════════════════════════════════════════════
clone(prs, orig_thanks)

# ════════════════════════════════════════════════════════════
# Remove original 15 slides (now at indices 0–14)
# ════════════════════════════════════════════════════════════
for _ in range(ORIG):
    del_slide(prs, 0)

prs.save(OUTPUT)
total = len(prs.slides)
print(f'Done. Saved to {OUTPUT}  ({total} slides)')
